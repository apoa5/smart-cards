import PyPDF2
from docx import Document
from pptx import Presentation
from .preprocess_text import preprocess_text

def extract_text_from_file(file, max_words=2000):
    filename = file.filename.lower()
    raw_text = ""

    if filename.endswith(".pdf"):
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                raw_text += page_text + "\n"

    elif filename.endswith(".txt"):
        raw_text = file.read().decode("utf-8")

    elif filename.endswith(".docx"):
        doc = Document(file)
        for para in doc.paragraphs:
            raw_text += para.text + "\n"

    elif filename.endswith(".pptx"):
        pres = Presentation(file)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    raw_text += shape.text + "\n"

    else:
        return ""  # unsupported file type

    return preprocess_text(raw_text, max_words=max_words)
